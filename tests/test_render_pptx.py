import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation

from pptx.util import Pt

from tools.render_pptx import render_pptx, SLIDE_WIDTH_PT, SLIDE_HEIGHT_PT, _add_placement
from tools.pipeline import MonthResult, SectionResult
from tools.blocks import grid_block, Placement

FIXED = {
    "title_lines": ["테스트 부제목", "테스트 배틀패스"],
    "period": "2026년 1월 1일 점검 후 ~ 2월 1일 AM 4:00",
    "grade_table": {"tiers": ["일반", "프리미엄", "로얄"], "prices": ["FREE", "24900캐쉬", "29900캐쉬"]},
    "notices": ["* 유의사항 1", "* 유의사항 2"],
}


def test_render_pptx_produces_openable_file(tmp_path):
    items = [{"name": f"아이템{i}", "image": None, "is_new": False} for i in range(6)]
    pages = grid_block(items, columns=3, icon_size=24.0)
    section = SectionResult(title="테스트 섹션", block_type="grid", item_count=6,
                             matched_image_count=0, text_only_count=6, pages=pages)
    result = MonthResult(month="999999", sections=[section])

    out_path = str(tmp_path / "test_output.pptx")
    render_pptx(FIXED, result, out_path)

    assert os.path.exists(out_path)
    prs = Presentation(out_path)
    from pptx.util import Pt
    assert prs.slide_width == Pt(SLIDE_WIDTH_PT)
    assert prs.slide_height == Pt(SLIDE_HEIGHT_PT)
    # 1 fixed-fields slide + 1 content slide (6 items fit on one grid page) + notices slide
    assert len(list(prs.slides)) == 3


def test_render_pptx_places_real_text_content(tmp_path):
    items = [{"name": "고유아이템이름123", "image": None, "is_new": False}]
    pages = grid_block(items, columns=1, icon_size=24.0)
    section = SectionResult(title="섹션제목", block_type="grid", item_count=1,
                             matched_image_count=0, text_only_count=1, pages=pages)
    result = MonthResult(month="999998", sections=[section])

    out_path = str(tmp_path / "test_output2.pptx")
    render_pptx(FIXED, result, out_path)
    prs = Presentation(out_path)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = " | ".join(all_text)
    assert "고유아이템이름123" in joined
    assert "섹션제목" in joined
    assert "테스트 배틀패스" in joined
    assert "유의사항 1" in joined


def test_render_pptx_inserts_new_badge_and_skips_render_error_sections(tmp_path):
    items = [{"name": "NEW아이템", "image": None, "is_new": True}]
    items += [{"name": f"일반아이템{i}", "image": None, "is_new": False} for i in range(3)]
    pages = grid_block(items, columns=3, icon_size=24.0)
    good_section = SectionResult(title="정상섹션", block_type="new_highlight", item_count=4,
                                  matched_image_count=0, text_only_count=4, pages=pages)
    bad_section = SectionResult(title="에러섹션", block_type="paired_columns", item_count=1,
                                 matched_image_count=0, text_only_count=1, pages=[],
                                 render_error="wrong item count")
    result = MonthResult(month="999997", sections=[good_section, bad_section])

    out_path = str(tmp_path / "test_output3.pptx")
    render_pptx(FIXED, result, out_path)
    prs = Presentation(out_path)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = " | ".join(all_text)
    assert "NEW아이템" in joined
    assert "에러섹션" not in joined  # render_error section must be skipped, not crash


def test_render_pptx_shows_pending_image_note_for_new_items_without_image(tmp_path):
    items = [{"name": "새아이템", "image": None, "is_new": True},
             {"name": "일반아이템", "image": None, "is_new": False}]
    pages = grid_block(items, columns=2, icon_size=24.0)
    section = SectionResult(title="섹션", block_type="grid", item_count=2,
                             matched_image_count=0, text_only_count=2, pages=pages)
    result = MonthResult(month="999996", sections=[section])

    out_path = str(tmp_path / "test_output4.pptx")
    render_pptx(FIXED, result, out_path)
    prs = Presentation(out_path)

    all_text = [s.text_frame.text for slide in prs.slides for s in slide.shapes if s.has_text_frame]
    joined = " | ".join(all_text)
    assert "이미지\n추후 전달 예정" in joined
    # only the new item without an image gets the note, not the ordinary one
    assert joined.count("이미지\n추후 전달 예정") == 1


def test_multiline_text_placement_styles_every_paragraph(tmp_path):
    # real bug: a "\n"-joined footnote ("* 문장1.\n* 문장2.") relied on
    # python-pptx auto-splitting text_frame.text on "\n" into paragraphs,
    # but only paragraphs[0] ever got an explicit font size — the
    # auto-created second paragraph fell back to the theme's default
    # (much larger) size and visually overflowed into the fixed content
    # below it. Every paragraph must get the same explicit styling.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    placement = Placement("text", 10.0, 10.0, 400.0, 40.0, "* 첫 번째 문장입니다.\n* 두 번째 문장입니다.")
    _add_placement(slide, placement)

    # "text" placements now also draw a card frame (see the test below) —
    # that's an autoshape, which python-pptx gives an (empty) text_frame
    # too, so skip to the shape that actually has the multiline text.
    tb = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text)
    tf = tb.text_frame
    assert len(tf.paragraphs) == 2
    assert tf.paragraphs[0].text == "* 첫 번째 문장입니다."
    assert tf.paragraphs[1].text == "* 두 번째 문장입니다."
    assert tf.paragraphs[0].font.size == Pt(9)
    assert tf.paragraphs[1].font.size == Pt(9)


def test_text_kind_placement_gets_its_own_card_frame(tmp_path):
    # real user comparison (2026-08-11): a text-only reward item (no source
    # image) rendered as bare floating text looked unfinished next to icon
    # items that all sit on a card — every item should read as "a card",
    # image or not.
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_placement(slide, Placement("text", 10.0, 10.0, 120.0, 40.0, "수묵화 대미지 스킨"))

    shapes = list(slide.shapes)
    assert len(shapes) == 2
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shapes[1].text_frame.text == "수묵화 대미지 스킨"


def test_caption_kind_placement_stays_unframed(tmp_path):
    # "caption" sits under an icon that already drew its own frame — must
    # NOT get a second frame of its own, unlike "text" above.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_placement(slide, Placement("caption", 10.0, 10.0, 120.0, 20.0, "아이템 이름"))

    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert shapes[0].text_frame.text == "아이템 이름"


def test_frame_placement_draws_card_with_no_picture_or_text(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_placement(slide, Placement("frame", 10.0, 10.0, 200.0, 100.0))
    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert not shapes[0].has_text_frame or shapes[0].text_frame.text == ""


def test_no_frame_meta_skips_card_border_on_image_placement(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_placement(slide, Placement("image", 10.0, 10.0, 80.0, 80.0, {"image": None}, meta={"no_frame": True}))
    assert len(list(slide.shapes)) == 0
