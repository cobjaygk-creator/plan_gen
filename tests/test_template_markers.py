import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from tools.template_markers import find_markers, remove_marker

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "samples", "template.pptx")


def test_finds_known_markers_on_real_template():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)

    red6 = [m for m in find_markers(slides[5]) if m.color == "red"]
    blue7 = [m for m in find_markers(slides[6]) if m.color == "blue"]
    blue8 = [m for m in find_markers(slides[7]) if m.color == "blue"]
    slide9_markers = find_markers(slides[8])

    assert len(red6) == 2
    assert len(blue7) == 1
    assert len(blue8) == 1
    assert sorted(m.color for m in slide9_markers) == ["blue", "red"]


def test_no_markers_on_boilerplate_slides():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)
    for idx in (0, 1, 2, 3, 4, 9):  # slides 1-5, 10 (0-indexed)
        assert find_markers(slides[idx]) == []


def test_remove_marker_deletes_the_shape(tmp_path):
    prs = Presentation(TEMPLATE_PATH)
    slide7 = list(prs.slides)[6]
    before = len(slide7.shapes)
    markers = find_markers(slide7)
    remove_marker(markers[0])
    assert len(slide7.shapes) == before - 1
    assert find_markers(slide7) == []


def test_removing_marker_does_not_touch_other_shapes(tmp_path):
    prs = Presentation(TEMPLATE_PATH)
    slide6 = list(prs.slides)[5]
    other_shapes_before = [
        (s.left, s.top, s.width, s.height)
        for s in slide6.shapes if find_markers_color(s) is None
    ]
    for m in find_markers(slide6):
        remove_marker(m)
    other_shapes_after = [(s.left, s.top, s.width, s.height) for s in slide6.shapes]
    for box in other_shapes_before:
        assert box in other_shapes_after


def find_markers_color(shape):
    from tools.template_markers import _marker_color
    return _marker_color(shape)
