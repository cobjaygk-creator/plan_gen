import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation

from tools.parse_fixed_fields import parse_fixed_fields
from tools.template_markers import find_markers, remove_marker
from tools.render_from_template import (
    apply_text_substitution, _shapes_under_marker, _overlap_ratio, render_from_template,
)
from tools.pipeline import MonthResult, SectionResult
from tools.blocks import grid_block

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "samples", "template.pptx")
REQUEST_605 = os.path.join(os.path.dirname(__file__), "..", "samples", "202605_request.xlsx")
REQUEST_512 = os.path.join(os.path.dirname(__file__), "..", "samples", "202512_request.xlsx")


def _slide6_and_9(prs):
    slides = list(prs.slides)
    return slides[5], slides[8]


def test_marker_never_matches_itself_as_a_candidate():
    # Regression: `shape is marker.shape` is unreliable across separate
    # .shapes iterations in python-pptx (fresh wrapper objects each time),
    # so the marker (empty text, has_text_frame=True) used to slip into
    # its own candidate list and steal a zip() slot from the real target.
    prs = Presentation(TEMPLATE_PATH)
    slide6, _ = _slide6_and_9(prs)
    markers = [m for m in find_markers(slide6) if m.color == "red"]
    for m in markers:
        targets = _shapes_under_marker(slide6, m)
        assert all(t.text_frame.text.strip() != "" for t in targets), \
            "marker's own empty-text shape leaked into its candidate list"


def test_text_substitution_replaces_title_period_and_notices():
    fixed = parse_fixed_fields(REQUEST_605)
    prs = Presentation(TEMPLATE_PATH)
    slide6, slide9 = _slide6_and_9(prs)

    for slide in (slide6, slide9):
        red = [m for m in find_markers(slide) if m.color == "red"]
        apply_text_substitution(slide, fixed, red)
        for m in red:
            remove_marker(m)

    slide6_text = " ".join(s.text_frame.text for s in slide6.shapes if s.has_text_frame)
    slide9_text = " ".join(s.text_frame.text for s in slide9.shapes if s.has_text_frame)

    assert fixed["title_lines"][0] in slide6_text
    assert fixed["title_lines"][1] in slide6_text
    assert "2026" in slide6_text and "5월" in slide6_text  # period substituted
    assert "시계태엽" not in slide6_text  # old December placeholder gone
    assert "6월 24일" in slide9_text  # notices end-date matches new period
    assert "유의사항" in slide9_text  # label itself must survive untouched


def test_substitution_does_not_corrupt_unrelated_label():
    # The "유의사항" label sits just above the notices marker and grazes
    # its edge — must never be treated as a substitution target.
    fixed = parse_fixed_fields(REQUEST_605)
    prs = Presentation(TEMPLATE_PATH)
    _, slide9 = _slide6_and_9(prs)
    red = [m for m in find_markers(slide9) if m.color == "red"]
    apply_text_substitution(slide9, fixed, red)
    labels = [s.text_frame.text for s in slide9.shapes if s.has_text_frame and s.text_frame.text.strip() == "유의사항"]
    assert labels == ["유의사항"]


def make_month_result(month, sections):
    return MonthResult(month=month, sections=sections)


def _grid_section(title, n_items, footnote=None):
    items = [{"name": f"{title}아이템{i}", "image": None, "is_new": False} for i in range(n_items)]
    pages = grid_block(items, columns=3, icon_size=24.0)
    return SectionResult(title=title, block_type="grid", item_count=n_items,
                          matched_image_count=0, text_only_count=n_items, pages=pages, footnote=footnote)


def test_render_from_template_fits_within_existing_slots(tmp_path):
    fixed = parse_fixed_fields(REQUEST_512)
    result = make_month_result("999900", [
        _grid_section("섹션A", 3),
        _grid_section("섹션B", 4),
    ])
    out_path = str(tmp_path / "out.pptx")
    render_from_template(fixed, result, TEMPLATE_PATH, out_path)

    prs = Presentation(out_path)
    # 2 pages: 1 goes to the reward-only slot (slide 7), 1 to the always-
    # last notices slide (slide 9) -> slide 8's reward-only slot goes
    # unused and gets removed entirely: 10 - 1 = 9
    assert len(list(prs.slides)) == 9

    all_text = " ".join(s.text_frame.text for slide in prs.slides for s in slide.shapes if s.has_text_frame)
    assert "섹션A" in all_text
    assert "섹션B" in all_text
    assert "유의사항" in all_text  # last slide's fixed content survived


def test_single_page_content_fits_the_actual_landing_slot_not_the_first_slot(tmp_path):
    # Real bug from an actual generated file: layout was computed against
    # the *first* blue slot's box (slides 7/8, ~460pt tall). A 15-item
    # grid fit fine there as a single page — but a lone page always lands
    # on the *last* slot (the notices slide), whose box is smaller
    # (~340pt tall) because it also carries the fixed 최종산출물/유의사항
    # content below it. The grid rendered right through that fixed text.
    fixed = parse_fixed_fields(REQUEST_512)
    result = make_month_result("999904", [_grid_section("배틀패스 신규 의상", 15)])
    out_path = str(tmp_path / "out.pptx")
    render_from_template(fixed, result, TEMPLATE_PATH, out_path)

    prs = Presentation(out_path)
    target_slide = next(
        slide for slide in prs.slides
        if any(s.has_text_frame and "유의사항" == s.text_frame.text.strip() for s in slide.shapes)
    )
    fixed_content_top = min(
        s.top / 12700 for s in target_slide.shapes
        if s.has_text_frame and s.text_frame.text.strip() == "유의사항"
    )
    item_caption_bottoms = [
        (s.top + s.height) / 12700 for s in target_slide.shapes
        if s.has_text_frame and s.text_frame.text.strip().startswith("배틀패스 신규 의상아이템")
    ]
    assert len(item_caption_bottoms) == 15  # all 15 captions actually rendered
    # every item caption must end above where the fixed 유의사항 label
    # begins — no overlap with the template's own fixed content
    assert all(bottom <= fixed_content_top + 1 for bottom in item_caption_bottoms), (
        max(item_caption_bottoms), fixed_content_top
    )


def test_render_from_template_clones_slides_when_content_overflows(tmp_path):
    fixed = parse_fixed_fields(REQUEST_512)
    result = make_month_result("999901", [
        _grid_section(f"섹션{i}", 3) for i in range(5)  # 5 pages, only 3 template slots
    ])
    out_path = str(tmp_path / "out.pptx")
    render_from_template(fixed, result, TEMPLATE_PATH, out_path)

    prs = Presentation(out_path)
    assert len(list(prs.slides)) == 12  # 10 + 2 cloned

    all_text = " ".join(s.text_frame.text for slide in prs.slides for s in slide.shapes if s.has_text_frame)
    for i in range(5):
        assert f"섹션{i}" in all_text
    assert "유의사항" in all_text
    # the notices slide must still be the second-to-last slide (the very
    # last slide is the template's own blank "마침"/ending slide, which
    # this module never touches or reorders relative to)
    second_to_last = list(prs.slides)[-2]
    text = " ".join(s.text_frame.text for s in second_to_last.shapes if s.has_text_frame)
    assert "유의사항" in text and "섹션4" in text


def test_render_from_template_removes_unused_slots(tmp_path):
    fixed = parse_fixed_fields(REQUEST_512)
    result = make_month_result("999902", [_grid_section("유일섹션", 2)])  # 1 page, 3 slots available
    out_path = str(tmp_path / "out.pptx")
    render_from_template(fixed, result, TEMPLATE_PATH, out_path)

    prs = Presentation(out_path)
    # 2 unused reward-only slots removed entirely -> 10 - 2 = 8 slides
    assert len(list(prs.slides)) == 8
    all_text = " ".join(s.text_frame.text for slide in prs.slides for s in slide.shapes if s.has_text_frame)
    assert "유일섹션" in all_text
    assert "유의사항" in all_text


def test_render_from_template_no_markers_remain(tmp_path):
    fixed = parse_fixed_fields(REQUEST_512)
    result = make_month_result("999903", [_grid_section("섹션", 3)])
    out_path = str(tmp_path / "out.pptx")
    render_from_template(fixed, result, TEMPLATE_PATH, out_path)

    prs = Presentation(out_path)
    for slide in prs.slides:
        assert find_markers(slide) == []
