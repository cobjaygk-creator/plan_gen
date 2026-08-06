"""Step 2 support: scan every sample result.pptx and pull out structural
signals (layout name, picture count/sizes, table count, key text snippets)
per slide, to support block-type classification without re-reading raw
pptx files by hand each time.
"""
import json
import glob
import os
from pptx import Presentation
from pptx.util import Emu

EMU_PER_PT = 12700


def collect_pictures(shapes):
    pics = []
    for s in shapes:
        if s.shape_type == 13:  # PICTURE
            pics.append((round(s.width / EMU_PER_PT, 1), round(s.height / EMU_PER_PT, 1)))
        if s.shape_type == 6:  # GROUP
            pics.extend(collect_pictures(s.shapes))
    return pics


def collect_text(shapes, limit=12):
    out = []
    for s in shapes:
        if s.has_text_frame:
            t = s.text_frame.text.strip()
            if t:
                out.append(t.replace("\n", " ")[:40])
        if s.shape_type == 6:
            out.extend(collect_text(s.shapes, limit))
    return out[:limit]


def collect_tables(shapes):
    n = 0
    for s in shapes:
        if s.has_table:
            n += 1
        if s.shape_type == 6:
            n += collect_tables(s.shapes)
    return n


def scan(path):
    prs = Presentation(path)
    slides_out = []
    for idx, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout.name
        pics = collect_pictures(slide.shapes)
        tables = collect_tables(slide.shapes)
        texts = collect_text(slide.shapes)
        has_new = any("NEW" in t.upper() for t in texts)
        slides_out.append({
            "slide_no": idx,
            "layout": layout,
            "picture_count": len(pics),
            "picture_sizes_pt": pics,
            "table_count": tables,
            "has_NEW_badge": has_new,
            "text_snippets": texts,
        })
    return slides_out


if __name__ == "__main__":
    result = {}
    for path in sorted(glob.glob("samples/*_result.pptx")):
        if path.startswith("samples/~$"):
            continue
        month = os.path.basename(path).split("_")[0]
        result[month] = scan(path)
    with open("templates/sample_scan.json", "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"scanned {len(result)} months -> templates/sample_scan.json")
