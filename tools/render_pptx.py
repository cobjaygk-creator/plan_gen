"""Final step: write the block engine's Placement coordinates (step 3) and
the fixed fields (step 5) out as a real .pptx file. Everything upstream of
this module only ever computed layout math or classified text — this is
the one place that actually touches python-pptx shape-creation calls.

Style values (colors/fonts/card frame/header bar) come from
templates/master_template.md's "렌더러용 확정 스타일 값" table — those were
reverse-engineered from the actual slide XML of samples/202605_result.pptx
because most of them are theme scheme-color references (schemeClr, not a
literal RGB) that python-pptx's high-level API can't read directly. Still
doesn't chase the master template's exact XML down to borders-on-every-
shape level — this covers what's visually load-bearing (card frame,
header bar, caption color, NEW styling), not full pixel parity.
"""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from tools.blocks.geometry import Placement, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_BOTTOM

SLIDE_WIDTH_PT = 960
SLIDE_HEIGHT_PT = 540
BODY_FONT = "맑은 고딕"

# templates/master_template.md 렌더러용 확정 스타일 값 표 그대로
CARD_FILL = RGBColor(0xFF, 0xFF, 0xFF)         # bg1
CARD_BORDER = RGBColor(0xBF, 0xBF, 0xBF)       # bg1 lumMod75% = "Background 1, Darker 25%"
CAPTION_COLOR = RGBColor(0x26, 0x26, 0x26)     # tx1 lumMod85%/lumOff15% = "Text 1, Lighter 15%"
HEADER_BG = RGBColor(0x40, 0x40, 0x40)         # tx1 lumMod75%/lumOff25% = "Text 1, Lighter 25%"
HEADER_TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # bg1
NEW_BADGE_COLOR = RGBColor(0xFF, 0x00, 0x00)
CARD_CORNER_RADIUS = 0.095  # adj≈9481 in prstGeom roundRect terms
CONTENT_BG_FILL = RGBColor(0xF2, 0xF2, 0xF2)   # bg1 lumMod95% = "Background 1, Darker 5%"
CONTENT_BG_BORDER = RGBColor(0xBF, 0xBF, 0xBF)  # same swatch as card border


def _add_content_background(slide):
    """Large rounded panel behind the whole content grid — source pptx's
    "모서리가 둥근 직사각형 89" shape (round2SameRect, bg1 lumMod95%/lumMod75%
    border). Was missing entirely before; every item just floated on the
    plain white slide with no panel behind it."""
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Pt(CONTENT_LEFT), Pt(CONTENT_TOP), Pt(CONTENT_WIDTH), Pt(CONTENT_BOTTOM - CONTENT_TOP),
    )
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = CONTENT_BG_FILL
    panel.line.color.rgb = CONTENT_BG_BORDER
    panel.line.width = Pt(0.25)
    panel.shadow.inherit = False
    return panel


def _item_name(ref) -> str:
    if isinstance(ref, dict):
        return ref.get("name", "")
    return str(ref) if ref is not None else ""


def _add_card_frame(slide, left, top, width, height):
    """White rounded-rect card frame — drawn behind every icon/image slot
    regardless of whether a real picture ends up on top of it, matching
    the source pptx (the card frame shape exists independently of 그림 N)."""
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.adjustments[0] = CARD_CORNER_RADIUS
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD_FILL
    frame.line.color.rgb = CARD_BORDER
    frame.line.width = Pt(0.25)
    frame.shadow.inherit = False
    return frame


def _add_placement(slide, placement: Placement):
    left, top = Pt(placement.left), Pt(placement.top)
    width, height = Pt(placement.width), Pt(placement.height)

    if placement.kind in ("icon", "image"):
        _add_card_frame(slide, left, top, width, height)
        image_path = placement.ref.get("image") if isinstance(placement.ref, dict) else None
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, left, top, width, height)
            return
        if isinstance(placement.ref, dict) and placement.ref.get("is_new"):
            # genuinely new item, art not ready yet — say so instead of a
            # silent blank box (matches source data's own "이미지 추후 전달 예정")
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = "이미지 추후 전달 예정"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(6)
            p.font.name = BODY_FONT
            p.font.color.rgb = CAPTION_COLOR
        # ordinary item with no matched image: leave the card frame empty —
        # the block engine always places a separate caption below/beside it
        # carrying the name (source pptx never puts text inside the card
        # shape itself either)
        return

    if placement.kind == "badge":
        # source pptx: no shape, just bold red rotated text (see master_template.md)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.rotation = -18
        tf = tb.text_frame
        tf.text = _item_name(placement.ref) or "NEW"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = NEW_BADGE_COLOR
        p.font.name = BODY_FONT
        return

    # caption / text
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = _item_name(placement.ref)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8)
    p.font.name = BODY_FONT
    p.font.color.rgb = CAPTION_COLOR


def _add_section_header(slide, text: str, left=Pt(20), top=Pt(20), width=Pt(500), height=Pt(36)):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = BODY_FONT
    p.font.color.rgb = HEADER_TEXT_COLOR


def _add_fixed_fields_slide(prs, blank_layout, fixed: dict):
    slide = prs.slides.add_slide(blank_layout)
    y = Pt(20)
    for line in fixed["title_lines"]:
        tb = slide.shapes.add_textbox(Pt(20), y, Pt(600), Pt(28))
        tb.text_frame.text = line
        tb.text_frame.paragraphs[0].font.size = Pt(20)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.name = BODY_FONT
        y += Pt(30)

    tb = slide.shapes.add_textbox(Pt(20), y, Pt(600), Pt(20))
    tb.text_frame.text = fixed["period"]
    tb.text_frame.paragraphs[0].font.size = Pt(11)
    tb.text_frame.paragraphs[0].font.name = BODY_FONT
    y += Pt(30)

    grade = fixed["grade_table"]
    rows, cols = 2, len(grade["tiers"])
    table_shape = slide.shapes.add_table(rows, cols, Pt(20), y, Pt(400), Pt(50))
    table = table_shape.table
    for c, tier in enumerate(grade["tiers"]):
        table.cell(0, c).text = tier
        table.cell(1, c).text = grade["prices"][c]


def _add_notices_slide(prs, blank_layout, fixed: dict):
    if not fixed["notices"]:
        return
    slide = prs.slides.add_slide(blank_layout)
    _add_section_header(slide, "유의사항")
    y = Pt(60)
    for line in fixed["notices"]:
        tb = slide.shapes.add_textbox(Pt(20), y, Pt(880), Pt(20))
        tb.text_frame.text = line
        tb.text_frame.paragraphs[0].font.size = Pt(10)
        tb.text_frame.paragraphs[0].font.name = BODY_FONT
        y += Pt(22)


def render_pptx(fixed: dict, month_result, out_path: str) -> str:
    """fixed: parse_fixed_fields() output. month_result: pipeline.MonthResult
    (sections with .pages already computed — this function does no
    classification/matching itself, just draws what it's given)."""
    prs = Presentation()
    prs.slide_width = Pt(SLIDE_WIDTH_PT)
    prs.slide_height = Pt(SLIDE_HEIGHT_PT)
    blank_layout = prs.slide_layouts[6]

    _add_fixed_fields_slide(prs, blank_layout, fixed)

    for section in month_result.sections:
        if section.render_error:
            continue
        for page in section.pages:
            slide = prs.slides.add_slide(blank_layout)
            _add_content_background(slide)
            _add_section_header(slide, section.title)
            for placement in page:
                _add_placement(slide, placement)

    _add_notices_slide(prs, blank_layout, fixed)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path
