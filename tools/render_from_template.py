"""Template-based renderer: works directly on samples/template.pptx (dashed
marker rectangles — red FF0000=text substitution, blue 00B0F0=reward
content area) instead of building slides from scratch. Everything else in
the template (header bar, background panel, right-side QA table, footer,
슬라이드 6's pricing table) is left completely untouched, per the user's
explicit requirement.

Pipeline: parse_fixed_fields + classify_month + locate_items + match_images
are unchanged (reused as-is) — only the rendering target changed, from
"build new slides" to "fill marked regions in the real template".
"""
import os
import re
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

from tools.template_markers import find_markers, remove_marker, Marker
from tools.blocks.dynamic_grid import fit_grid_pages
from tools.blocks.text_list import text_list_block
from tools.blocks.few_preview import few_preview_block
from tools.blocks.paired_columns import paired_columns_block
from tools.render_pptx import _add_placement, _add_section_header

EMU_PER_PT = 12700
HEADER_HEIGHT = 36.0
HEADER_GAP = 12.0

PERIOD_RE = re.compile(r".*점검.*~.*(AM|PM|오전|오후).*\d{1,2}:\d{2}.*")
NOTICE_START_RE = re.compile(r"^\s*\*")


# ---------------------------------------------------------------- text ----

def _bbox_pt(shape):
    return (shape.left / EMU_PER_PT, shape.top / EMU_PER_PT, shape.width / EMU_PER_PT, shape.height / EMU_PER_PT)


def _overlap_ratio(shape_box, marker_box) -> float:
    """Fraction of shape_box's own area that falls inside marker_box."""
    sx, sy, sw, sh = shape_box
    mx, my, mw, mh = marker_box
    ox = max(0.0, min(sx + sw, mx + mw) - max(sx, mx))
    oy = max(0.0, min(sy + sh, my + mh) - max(sy, my))
    shape_area = sw * sh
    return (ox * oy) / shape_area if shape_area else 0.0


def _shapes_under_marker(slide, marker: Marker, min_overlap_ratio: float = 0.4):
    """Requires most of the candidate shape's own area to fall inside the
    marker, not just any pixel of overlap — a shape that merely grazes the
    marker's edge (e.g. a label sitting just above it, whose bottom edge
    dips a few points into the marker's top) must NOT count. Found via a
    real case: the "유의사항" label above the slide-9 notices marker
    overlapped it by a few points and got its text overwritten instead of
    the actual notices paragraph underneath."""
    marker_box = (marker.left, marker.top, marker.width, marker.height)
    found = []
    for shape in slide.shapes:
        # `shape is marker.shape` is NOT reliable here: python-pptx hands
        # back a fresh wrapper object on every `.shapes` iteration, so two
        # wrappers around the very same XML element still fail an `is`
        # check. Compare the underlying element instead. (Found via a real
        # bug: the marker itself — empty text, but has_text_frame=True —
        # was slipping into its own candidate list and stealing a zip()
        # slot meant for the real title/period shape.)
        if shape._element is marker.shape._element or shape.left is None or not shape.has_text_frame:
            continue
        box = _bbox_pt(shape)
        if _overlap_ratio(box, marker_box) >= min_overlap_ratio:
            found.append(shape)
    found.sort(key=lambda s: s.top)
    return found


def _replace_single_line(shape, text: str) -> None:
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in list(p0.runs[1:]):
            r._r.getparent().remove(r._r)
    else:
        p0.text = text
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def _replace_multiline(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    template_run = p0.runs[0] if p0.runs else None
    font_name = template_run.font.name if template_run else None
    font_size = template_run.font.size if template_run else None
    font_bold = template_run.font.bold if template_run else None

    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    lines = lines or [""]
    _replace_single_line(shape, lines[0])
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if font_name:
            r.font.name = font_name
        if font_size:
            r.font.size = font_size
        if font_bold is not None:
            r.font.bold = font_bold


def apply_text_substitution(slide, fixed: dict, red_markers: list[Marker]) -> None:
    """Classifies each red marker's role by the *existing* text underneath
    it (this template.pptx already carries real December-2025 placeholder
    content, not lorem-ipsum, so pattern-matching it is reliable) rather
    than assuming a fixed marker order/position, which would break if the
    template's layout ever shifts slightly."""
    for marker in red_markers:
        targets = _shapes_under_marker(slide, marker)
        if not targets:
            continue
        combined = " ".join(t.text_frame.text for t in targets).strip()
        if PERIOD_RE.match(combined):
            _replace_single_line(targets[0], fixed["period"])
        elif NOTICE_START_RE.match(combined) or len(combined) > 30:
            _replace_multiline(targets[0], fixed["notices"])
        else:
            for shape, text in zip(targets, fixed["title_lines"]):
                _replace_single_line(shape, text)


# ------------------------------------------------------------ layout ----

def _content_box(box):
    left, top, width, height = box
    content_top = top + HEADER_HEIGHT + HEADER_GAP
    return (left, content_top, width, height - HEADER_HEIGHT - HEADER_GAP)


def _render_section_pages(section, items: list[dict], box):
    content_box = _content_box(box)
    bt = section.block_type
    if bt in ("grid", "new_highlight"):
        return fit_grid_pages(items, *content_box)
    if bt == "text_list":
        return text_list_block(items, columns=3)
    if bt == "few_preview":
        return few_preview_block(items, box=content_box)
    if bt == "paired_columns":
        pair_items, sub_items = items[:2], items[2:] or None
        return paired_columns_block(pair_items, sub_items=sub_items, footnote=section.footnote, box=content_box)
    raise ValueError(f"unknown block_type: {bt!r}")


def _items_from_pages(section) -> list[dict]:
    """section.pages already has items baked into Placement.ref (from
    pipeline.py's _items_with_images) — pull the unique item dicts back
    out in original order rather than re-deriving them, so this module
    doesn't need its own copy of the classify/match logic."""
    seen = set()
    items = []
    for page in section.pages:
        for placement in page:
            if placement.kind in ("icon", "image") and isinstance(placement.ref, dict):
                key = placement.ref.get("name")
                if key not in seen:
                    seen.add(key)
                    items.append(placement.ref)
    return items


# -------------------------------------------------------- slide surgery ----

def _duplicate_slide(prs, source_slide):
    """Clones a slide's full shape tree onto a brand-new slide (appended
    at the end of the deck by add_slide — caller is responsible for
    reordering it into place)."""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree)[2:]:  # keep nvGrpSpPr, grpSpPr; drop layout-inherited placeholders
        sp_tree.remove(child)
    for shape_el in list(source_slide.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(shape_el))
    return new_slide


def _move_slide_before(prs, slide_to_move, before_slide) -> None:
    xml_slides = prs.slides._sldIdLst
    move_el, before_el = None, None
    for sld in xml_slides:
        related = prs.part.related_part(sld.get(qn("r:id")))
        if related is slide_to_move.part:
            move_el = sld
        if related is before_slide.part:
            before_el = sld
    if move_el is not None and before_el is not None:
        xml_slides.remove(move_el)
        before_el.addprevious(move_el)


def _remove_slide(prs, slide) -> None:
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        if prs.part.related_part(sld.get(qn("r:id"))) is slide.part:
            xml_slides.remove(sld)
            break


# -------------------------------------------------------------- main ----

def render_from_template(fixed: dict, month_result, template_path: str, out_path: str) -> str:
    prs = Presentation(template_path)
    original_slides = list(prs.slides)

    for slide in original_slides:
        red = [m for m in find_markers(slide) if m.color == "red"]
        if red:
            apply_text_substitution(slide, fixed, red)
        for m in red:
            remove_marker(m)

    blue_slots = []  # (slide, Marker) in document order
    for slide in original_slides:
        for m in find_markers(slide):
            if m.color == "blue":
                blue_slots.append((slide, m))
    if not blue_slots:
        raise ValueError(f"{template_path}: no blue (00B0F0) content markers found")

    sections = [s for s in month_result.sections if not s.render_error]

    # Layout math has to use ONE reference box size for every page, since
    # which page lands on which slide isn't known until after layout runs.
    # Real bug found via actual output: using the *first* slot's box
    # (slides 7/8, ~460pt square) computed a single page for a 15-item
    # section that fit fine there — but that lone page is always the "last
    # page" (see below), which always lands on the notices slide's smaller
    # box (~340pt tall), so it overflowed straight through the fixed
    # "최종 산출물은..." button and 유의사항 text underneath. Using the
    # *smallest* box among all slots as the shared reference guarantees
    # whatever gets computed fits whichever slot it actually lands on —
    # slightly under-uses the larger slots, but under-using space is a far
    # smaller problem than overlapping fixed template content.
    smallest_marker = min((m for _, m in blue_slots), key=lambda m: m.width * m.height)
    box0 = (smallest_marker.left, smallest_marker.top, smallest_marker.width, smallest_marker.height)

    all_pages = []  # (section_title, placements) — laid out against box0
    for section in sections:
        items = _items_from_pages(section)
        for page in _render_section_pages(section, items, box0):
            all_pages.append((section.title, page))

    last_slide, last_marker = blue_slots[-1]  # always carries notices/final
    # content too (see template inspection) — this slide must stay in the
    # deck and must always receive the *last* reward page, never be treated
    # as just another interchangeable reward slot.
    reward_only_slots = list(blue_slots[:-1])

    if not all_pages:
        # no reward content at all — clear every blue marker, change nothing else
        for slide, marker in blue_slots:
            remove_marker(marker)
    else:
        last_page = all_pages[-1]
        other_pages = all_pages[:-1]

        if len(other_pages) <= len(reward_only_slots):
            used_slots = reward_only_slots[:len(other_pages)]
            unused_slots = reward_only_slots[len(other_pages):]
        else:
            used_slots = list(reward_only_slots)
            unused_slots = []
            clone_source_slide, _ = reward_only_slots[-1] if reward_only_slots else (last_slide, last_marker)
            extra_needed = len(other_pages) - len(reward_only_slots)
            for _ in range(extra_needed):
                new_slide = _duplicate_slide(prs, clone_source_slide)
                _move_slide_before(prs, new_slide, last_slide)
                new_marker = next(m for m in find_markers(new_slide) if m.color == "blue")
                used_slots.append((new_slide, new_marker))

        for (title, placements), (slide, marker) in zip(other_pages, used_slots):
            box = (marker.left, marker.top, marker.width, marker.height)
            _add_section_header(slide, title, left=Pt(box[0]), top=Pt(box[1]), width=Pt(box[2]), height=Pt(HEADER_HEIGHT))
            for placement in placements:
                _add_placement(slide, placement)
            remove_marker(marker)

        title, placements = last_page
        box = (last_marker.left, last_marker.top, last_marker.width, last_marker.height)
        _add_section_header(last_slide, title, left=Pt(box[0]), top=Pt(box[1]), width=Pt(box[2]), height=Pt(HEADER_HEIGHT))
        for placement in placements:
            _add_placement(last_slide, placement)
        remove_marker(last_marker)

        # reward-only slots that ended up unused (fewer other_pages than
        # slots) never had reward content or any other purpose — remove
        # the whole slide, not just its marker
        for slide, marker in unused_slots:
            _remove_slide(prs, slide)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path
