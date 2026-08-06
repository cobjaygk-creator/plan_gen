"""Step 1: extract master template values (color / font / coordinate) from a
known-good sample result.pptx. Outputs a structured JSON that later steps
(block engine, regression scorer) read instead of touching the pptx again.
"""
import json
import sys
from pptx import Presentation
from pptx.util import Emu


def rgb_of(fill):
    try:
        if fill.type is None:
            return None
        color = fill.fore_color
        if color.type is not None and hasattr(color, "rgb"):
            return str(color.rgb)
    except Exception:
        return None
    return None


def font_info(run):
    f = run.font
    color = None
    try:
        if f.color and f.color.type is not None:
            color = str(f.color.rgb)
    except Exception:
        pass
    return {
        "text_sample": run.text[:20],
        "name": f.name,
        "size_pt": f.size.pt if f.size else None,
        "bold": f.bold,
        "color": color,
    }


def dump_shape(shape, depth=0):
    entry = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left_emu": shape.left,
        "top_emu": shape.top,
        "width_emu": shape.width,
        "height_emu": shape.height,
    }
    if shape.has_text_frame:
        fonts = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    fonts.append(font_info(run))
        if fonts:
            entry["fonts"] = fonts[:3]
    try:
        if shape.fill and shape.fill.type is not None:
            entry["fill_rgb"] = rgb_of(shape.fill)
    except Exception:
        pass
    if shape.shape_type == 6:  # GROUP
        entry["children"] = [dump_shape(s, depth + 1) for s in shape.shapes]
    if shape.has_table:
        tbl = shape.table
        entry["table"] = {
            "rows": len(tbl.rows),
            "cols": len(tbl.columns),
            "col_widths_emu": [c.width for c in tbl.columns],
            "row_heights_emu": [r.height for r in tbl.rows],
        }
    return entry


def extract(path):
    prs = Presentation(path)
    out = {
        "source": path,
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": [],
    }
    for idx, slide in enumerate(prs.slides, start=1):
        slide_entry = {
            "slide_no": idx,
            "layout_name": slide.slide_layout.name,
            "shapes": [dump_shape(s) for s in slide.shapes],
        }
        out["slides"].append(slide_entry)
    return out


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "samples/202605_result.pptx"
    data = extract(path)
    out_path = "templates/master_template_raw.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"wrote {out_path}: {len(data['slides'])} slides")
