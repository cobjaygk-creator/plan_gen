"""Detect overlapping shapes on each slide of a pptx file.

v2 — the first version (see git history) flagged group-shape-vs-own-child
containment and background-panel-vs-content containment as "overlaps",
which is just normal nesting/layering, not a bug. This version:
  1. Only compares leaf shapes (flattens into group children, never
     compares a group's own bounding box against its children).
  2. Ignores full-containment pairs (one box entirely inside another —
     that's a background panel under content, by design).
  3. Only flags partial overlap above a minimum area, which is what an
     actual layout bug (two items crammed into the same space) looks like.

Used both to pick a bug-free master template sample (step 1) and as one of
the regression scoring signals (step 4).
"""
from pptx import Presentation

CONTAINMENT_TOLERANCE = 0  # pt-EMU slack when deciding "fully inside"


def leaf_boxes(shapes):
    boxes = []
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP — recurse, never compare the group itself
            boxes.extend(leaf_boxes(shape.shapes))
            continue
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            continue
        boxes.append((shape.shape_id, shape.name, shape.left, shape.top,
                      shape.left + shape.width, shape.top + shape.height))
    return boxes


def overlap_area(a, b):
    _, _, ax0, ay0, ax1, ay1 = a
    _, _, bx0, by0, bx1, by1 = b
    ox = max(0, min(ax1, bx1) - max(ax0, bx0))
    oy = max(0, min(ay1, by1) - max(ay0, by0))
    return ox * oy


def fully_contains(a, b):
    """True if box a fully contains box b (or vice versa is checked by caller)."""
    _, _, ax0, ay0, ax1, ay1 = a
    _, _, bx0, by0, bx1, by1 = b
    return (ax0 - CONTAINMENT_TOLERANCE <= bx0 and ay0 - CONTAINMENT_TOLERANCE <= by0 and
            ax1 + CONTAINMENT_TOLERANCE >= bx1 and ay1 + CONTAINMENT_TOLERANCE >= by1)


def find_overlaps(path, min_area_pt2=15.0):
    """min_area_pt2: ignore overlaps smaller than this (border/AA noise)."""
    min_area_emu = min_area_pt2 * (12700 ** 2)
    prs = Presentation(path)
    report = []
    for idx, slide in enumerate(prs.slides, start=1):
        boxes = leaf_boxes(slide.shapes)
        pairs = []
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                if fully_contains(a, b) or fully_contains(b, a):
                    continue
                area = overlap_area(a, b)
                if area > min_area_emu:
                    pairs.append((a[1], b[1], area))
        if pairs:
            report.append((idx, pairs))
    return report


if __name__ == "__main__":
    import sys
    for name in sys.argv[1:]:
        path = f"samples/{name}_result.pptx"
        report = find_overlaps(path)
        if not report:
            print(f"{name}: OK (no significant partial overlaps)")
        else:
            print(f"{name}: {len(report)} slide(s) with overlaps")
            for slide_no, pairs in report:
                for a, b, area in pairs:
                    pt2 = area / (12700 ** 2)
                    print(f"  slide {slide_no}: '{a}' x '{b}' overlap~{pt2:.1f}pt2")
